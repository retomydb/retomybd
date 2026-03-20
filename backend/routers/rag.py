"""
retomY — RAG (Retrieval Augmented Generation) Router
Handles file upload, text extraction, and context-aware chat with Ollama models.
Supports: PDF, DOCX, XLSX, PPTX, CSV, JSON, TXT, code files, images, audio.
Audio is transcribed via Ollama whisper or sent as description prompt.
"""

from fastapi import APIRouter, UploadFile, File, Form, HTTPException, status
from fastapi.responses import StreamingResponse, JSONResponse
from pydantic import BaseModel, Field
from typing import Optional, Any
from core.config import get_settings
import structlog
import httpx
import json
import io
import os
import tempfile
import base64
import mimetypes

logger = structlog.get_logger()
settings = get_settings()
router = APIRouter(prefix="/rag", tags=["RAG"])

OLLAMA_URL = settings.OLLAMA_BASE_URL

# Max file size: 50MB
MAX_FILE_SIZE = 50 * 1024 * 1024

# ── Shared HTTP client ────────────────────────────────────────────────────────

_http_client: httpx.AsyncClient | None = None


def _get_client() -> httpx.AsyncClient:
    global _http_client
    if _http_client is None or _http_client.is_closed:
        _http_client = httpx.AsyncClient(
            timeout=httpx.Timeout(300.0, connect=10.0),
            limits=httpx.Limits(max_connections=20, max_keepalive_connections=10),
            follow_redirects=True,
        )
    return _http_client


# ── Text extraction for various file types ────────────────────────────────────

def extract_text_from_pdf(data: bytes) -> str:
    """Extract text from PDF bytes."""
    from PyPDF2 import PdfReader
    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[Page {i+1}]\n{text.strip()}")
    return "\n\n".join(pages) if pages else "[PDF: No extractable text found]"


def extract_text_from_docx(data: bytes) -> str:
    """Extract text from DOCX bytes."""
    from docx import Document
    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract from tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    return "\n\n".join(paragraphs) if paragraphs else "[DOCX: No text found]"


def extract_text_from_xlsx(data: bytes) -> str:
    """Extract text from XLSX bytes."""
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows_text.append(" | ".join(cells))
        if rows_text:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows_text))
    wb.close()
    return "\n\n".join(sheets) if sheets else "[XLSX: No data found]"


def extract_text_from_pptx(data: bytes) -> str:
    """Extract text from PPTX bytes."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
    return "\n\n".join(slides) if slides else "[PPTX: No text found]"


def extract_text_from_csv(data: bytes) -> str:
    """Extract text from CSV bytes."""
    import csv
    import chardet
    detected = chardet.detect(data[:10000])
    encoding = detected.get("encoding", "utf-8") or "utf-8"
    text = data.decode(encoding, errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = []
    for i, row in enumerate(reader):
        if i > 500:  # Cap at 500 rows
            rows.append(f"... ({i} rows total, truncated)")
            break
        rows.append(" | ".join(row))
    return "\n".join(rows) if rows else "[CSV: Empty file]"


def extract_text_from_json(data: bytes) -> str:
    """Extract text from JSON bytes."""
    import chardet
    detected = chardet.detect(data[:10000])
    encoding = detected.get("encoding", "utf-8") or "utf-8"
    text = data.decode(encoding, errors="replace")
    try:
        obj = json.loads(text)
        formatted = json.dumps(obj, indent=2, ensure_ascii=False)
        if len(formatted) > 50000:
            formatted = formatted[:50000] + "\n... [truncated]"
        return formatted
    except json.JSONDecodeError:
        return text[:50000]


def extract_text_from_text(data: bytes) -> str:
    """Extract text from plain text / code files."""
    import chardet
    detected = chardet.detect(data[:10000])
    encoding = detected.get("encoding", "utf-8") or "utf-8"
    text = data.decode(encoding, errors="replace")
    if len(text) > 50000:
        text = text[:50000] + "\n... [truncated]"
    return text


# Supported file extensions mapped to extractors
EXTRACTORS = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".doc": extract_text_from_docx,
    ".xlsx": extract_text_from_xlsx,
    ".xls": extract_text_from_xlsx,
    ".pptx": extract_text_from_pptx,
    ".csv": extract_text_from_csv,
    ".json": extract_text_from_json,
    ".jsonl": extract_text_from_json,
    # Code / text files
    ".txt": extract_text_from_text,
    ".md": extract_text_from_text,
    ".py": extract_text_from_text,
    ".js": extract_text_from_text,
    ".ts": extract_text_from_text,
    ".tsx": extract_text_from_text,
    ".jsx": extract_text_from_text,
    ".html": extract_text_from_text,
    ".css": extract_text_from_text,
    ".sql": extract_text_from_text,
    ".yaml": extract_text_from_text,
    ".yml": extract_text_from_text,
    ".xml": extract_text_from_text,
    ".ini": extract_text_from_text,
    ".toml": extract_text_from_text,
    ".sh": extract_text_from_text,
    ".bash": extract_text_from_text,
    ".java": extract_text_from_text,
    ".c": extract_text_from_text,
    ".cpp": extract_text_from_text,
    ".h": extract_text_from_text,
    ".go": extract_text_from_text,
    ".rs": extract_text_from_text,
    ".rb": extract_text_from_text,
    ".php": extract_text_from_text,
    ".swift": extract_text_from_text,
    ".kt": extract_text_from_text,
    ".r": extract_text_from_text,
    ".scala": extract_text_from_text,
    ".lua": extract_text_from_text,
    ".log": extract_text_from_text,
    ".env": extract_text_from_text,
    ".cfg": extract_text_from_text,
    ".conf": extract_text_from_text,
    ".dockerfile": extract_text_from_text,
}

# Image extensions (these go to Ollama as base64 for vision models)
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".svg"}

# Audio extensions
AUDIO_EXTENSIONS = {".mp3", ".wav", ".ogg", ".m4a", ".flac", ".webm", ".aac", ".wma"}


def get_file_extension(filename: str) -> str:
    """Get normalized file extension."""
    _, ext = os.path.splitext(filename.lower())
    return ext


def extract_file_content(filename: str, data: bytes) -> dict:
    """
    Extract content from a file. Returns dict with:
    - type: 'text' | 'image' | 'audio' | 'unsupported'
    - content: extracted text or base64 data
    - filename: original filename
    """
    ext = get_file_extension(filename)

    if ext in IMAGE_EXTENSIONS:
        # Encode image as base64 for vision models
        b64 = base64.b64encode(data).decode("utf-8")
        mime = mimetypes.guess_type(filename)[0] or "image/png"
        return {
            "type": "image",
            "content": b64,
            "mime_type": mime,
            "filename": filename,
            "size": len(data),
        }

    if ext in AUDIO_EXTENSIONS:
        b64 = base64.b64encode(data).decode("utf-8")
        mime = mimetypes.guess_type(filename)[0] or "audio/wav"
        return {
            "type": "audio",
            "content": b64,
            "mime_type": mime,
            "filename": filename,
            "size": len(data),
        }

    extractor = EXTRACTORS.get(ext)
    if extractor:
        try:
            text = extractor(data)
            return {
                "type": "text",
                "content": text,
                "filename": filename,
                "size": len(data),
                "char_count": len(text),
            }
        except Exception as e:
            logger.error("file_extract_error", filename=filename, ext=ext, error=str(e))
            return {
                "type": "text",
                "content": f"[Error extracting {filename}: {str(e)}]",
                "filename": filename,
                "size": len(data),
            }

    # Fallback: try as text
    try:
        text = extract_text_from_text(data)
        return {
            "type": "text",
            "content": text,
            "filename": filename,
            "size": len(data),
            "char_count": len(text),
        }
    except Exception:
        return {"type": "unsupported", "filename": filename, "size": len(data)}


# ── Endpoints ─────────────────────────────────────────────────────────────────

@router.post("/upload")
async def upload_and_extract(file: UploadFile = File(...)):
    """
    Upload a file and extract its text content.
    Returns extracted text that the frontend can include as context in chat messages.
    """
    if not file.filename:
        raise HTTPException(400, "No filename provided")

    data = await file.read()
    if len(data) > MAX_FILE_SIZE:
        raise HTTPException(413, f"File too large. Max {MAX_FILE_SIZE // (1024*1024)}MB")

    result = extract_file_content(file.filename, data)

    if result["type"] == "unsupported":
        raise HTTPException(
            415,
            f"Unsupported file type: {get_file_extension(file.filename)}. "
            f"Supported: {', '.join(sorted(set(list(EXTRACTORS.keys()) + list(IMAGE_EXTENSIONS) + list(AUDIO_EXTENSIONS))))}"
        )

    return result


@router.post("/chat")
async def rag_chat(
    file: Optional[UploadFile] = File(None),
    model: str = Form(...),
    message: str = Form(""),
    context: str = Form(""),
    history: str = Form("[]"),
    stream: str = Form("true"),
    temperature: float = Form(0.7),
    top_p: float = Form(0.9),
    max_tokens: int = Form(2048),
):
    """
    RAG-enabled chat: send a message with optional file context to an Ollama model.
    The file is parsed server-side and its text is injected as context.
    """
    client = _get_client()

    # Parse history
    try:
        chat_history = json.loads(history)
    except json.JSONDecodeError:
        chat_history = []

    # Build the prompt with file context
    file_context = ""
    images_b64: list[str] = []

    # Handle uploaded file
    if file and file.filename:
        data = await file.read()
        if len(data) > MAX_FILE_SIZE:
            raise HTTPException(413, f"File too large. Max {MAX_FILE_SIZE // (1024*1024)}MB")

        extracted = extract_file_content(file.filename, data)

        if extracted["type"] == "text":
            file_context = (
                f"\n\n--- FILE: {extracted['filename']} ---\n"
                f"{extracted['content']}\n"
                f"--- END FILE ---\n"
            )
        elif extracted["type"] == "image":
            images_b64.append(extracted["content"])
            file_context = f"\n[Image attached: {extracted['filename']}]\n"
        elif extracted["type"] == "audio":
            # For audio, we'll transcribe using the model's understanding
            file_context = f"\n[Audio file attached: {extracted['filename']} ({extracted['size']} bytes)]\n"

    # Also include pre-extracted context from frontend
    if context:
        file_context = f"\n\n--- DOCUMENT CONTEXT ---\n{context}\n--- END CONTEXT ---\n" + file_context

    # Build messages
    messages = []

    # System message with RAG instructions
    system_content = (
        "You are a helpful AI assistant with RAG (Retrieval Augmented Generation) capabilities. "
        "When the user provides file content or document context, analyze it thoroughly and provide "
        "detailed, accurate responses based on the content. Reference specific parts of the document "
        "when answering. If the user asks about something not in the provided context, clearly state "
        "that the information is not in the provided document."
    )
    messages.append({"role": "system", "content": system_content})

    # Add chat history
    for msg in chat_history:
        if isinstance(msg, dict) and "role" in msg and "content" in msg:
            messages.append({"role": msg["role"], "content": msg["content"]})

    # Build user message with context
    user_content = message.strip()
    if file_context:
        user_content = file_context + "\n\n" + (user_content or "Please analyze this file and provide a detailed explanation.")
    elif not user_content:
        raise HTTPException(400, "No message or file provided")

    user_msg: dict[str, Any] = {"role": "user", "content": user_content}
    if images_b64:
        user_msg["images"] = images_b64
    messages.append(user_msg)

    # Build payload
    payload: dict[str, Any] = {
        "model": model,
        "messages": messages,
        "stream": stream.lower() == "true",
        "options": {
            "temperature": temperature,
            "top_p": top_p,
            "num_predict": max_tokens,
        },
    }

    try:
        if payload["stream"]:
            return StreamingResponse(
                _stream_ollama(client, f"{OLLAMA_URL}/api/chat", payload),
                media_type="text/event-stream",
            )
        else:
            resp = await client.post(f"{OLLAMA_URL}/api/chat", json=payload)
            resp.raise_for_status()
            data = resp.json()
            return {
                "message": data.get("message", {}),
                "model": data.get("model", model),
                "done": data.get("done", True),
                "total_duration": data.get("total_duration"),
                "eval_count": data.get("eval_count"),
            }
    except httpx.ConnectError:
        raise HTTPException(502, "Cannot reach Ollama server. Is the container running?")
    except httpx.HTTPStatusError as e:
        err_text = e.response.text[:500] if e.response else str(e)
        logger.error("rag_chat_error", model=model, error=err_text)
        raise HTTPException(502, f"Ollama error: {err_text}")


@router.post("/transcribe")
async def transcribe_audio(file: UploadFile = File(...), model: str = Form("llama3:8b")):
    """
    Transcribe audio by sending it to Ollama with a transcription prompt.
    Since Ollama doesn't have a native Whisper API, we use the browser's
    Web Speech API on the frontend. This endpoint is a fallback that describes
    the audio metadata.
    """
    if not file.filename:
        raise HTTPException(400, "No filename provided")

    data = await file.read()
    if len(data) > MAX_FILE_SIZE:
        raise HTTPException(413, "Audio file too large")

    ext = get_file_extension(file.filename)
    if ext not in AUDIO_EXTENSIONS:
        raise HTTPException(415, f"Not an audio file: {ext}")

    # Return metadata — actual transcription happens via Web Speech API on frontend
    return {
        "filename": file.filename,
        "size": len(data),
        "mime_type": mimetypes.guess_type(file.filename)[0] or "audio/wav",
        "message": "Audio transcription is handled client-side via Web Speech API. "
                   "Use the microphone button to record and transcribe speech directly.",
    }


@router.get("/supported-formats")
async def supported_formats():
    """List all supported file formats for upload."""
    return {
        "document": [".pdf", ".docx", ".doc", ".pptx"],
        "spreadsheet": [".xlsx", ".xls", ".csv"],
        "data": [".json", ".jsonl"],
        "code": [".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".css", ".sql",
                 ".java", ".c", ".cpp", ".go", ".rs", ".rb", ".php", ".swift",
                 ".kt", ".scala", ".lua", ".sh", ".bash"],
        "text": [".txt", ".md", ".yaml", ".yml", ".xml", ".toml", ".ini",
                 ".log", ".env", ".cfg", ".conf"],
        "image": sorted(IMAGE_EXTENSIONS),
        "audio": sorted(AUDIO_EXTENSIONS),
    }


# ── Helpers ───────────────────────────────────────────────────────────────────

async def _stream_ollama(client: httpx.AsyncClient, url: str, payload: dict):
    """Stream SSE-formatted chunks from Ollama."""
    payload["stream"] = True
    async with client.stream("POST", url, json=payload) as resp:
        resp.raise_for_status()
        async for line in resp.aiter_lines():
            if not line.strip():
                continue
            yield f"data: {line}\n\n"
    yield "data: [DONE]\n\n"
